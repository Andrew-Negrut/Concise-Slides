import openai
from pptx import Presentation
from environs import Env

def get_concise_paragraph(paragraph):
    env = Env()
    env.read_env()

    openai.api_key = env.str("OPENAI_API_KEY")

    prompt = """Shorten the following text to make it concise while retaining its meaning (the result should be no more than 6 words. The result should have the same meaning as the original text.
    Skip names, and leave single words or phrases as they are): """ + paragraph

    messages = [{"role": "user", "content": prompt}]
    response = openai.ChatCompletion.create(model="gpt-4", messages=messages)

    return response.choices[0].message["content"]    

def get_concise(file):

    prs = Presentation(file)
    for slide_number, slide in enumerate(prs.slides):
        if slide_number == 0: 
            continue
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if not paragraph.text.strip():
                        continue
                    if len(paragraph.text.strip()) < 20:
                        continue
                    print("before: ", paragraph.text)
                    paragraph.text = get_concise_paragraph(paragraph.text)
                    print("after: ", paragraph.text)      

    return prs